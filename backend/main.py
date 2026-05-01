import base64
import json
import os
import re
import tempfile
import uuid
from typing import Any

import fitz
from fastapi import Body, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from openai import AuthenticationError, OpenAI, OpenAIError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SYSTEM_PROMPT = """You are a scientific claim extraction engine.

You MUST produce structured scientific claims.

Rules:
- Never copy sentences from the paper
- Never output generic statements
- Always infer the mechanism if not explicit
- Prefer specific, mechanistic, causal claims
- Empty fields are forbidden
- Every claim must include paper-specific nouns: population/model, manipulation/condition, measured variable, or mechanism.
- If a sentence could fit many papers, rewrite it until it is specific to this paper.
- Reject vague verbs like "affects", "influences", "is associated with", or "may help" unless paired with a concrete measured effect.
- Prefer named variables, task names, cell types, brain regions, biomarkers, interventions, datasets, or parameters from the paper.

If unclear:
→ infer the most likely scientific interpretation

Reject generic outputs.
The final infographic must fit on one A4 page. Compress aggressively before adding detail.

Output format:

{
  "full_structured_claims": {
    "thesis": "",
    "why_it_matters": [],
    "study_design": {
      "model_system": [],
      "methods": [],
      "sample": "",
      "manipulation": "",
      "measures": []
    },
    "core_evidence": [],
    "mechanism": "",
    "boundary_conditions": [],
    "generalizable_insight": ""
  },
  "infographic_claims": {
    "thesis": "",
    "why_it_matters": [],
    "study_design": {
      "model_system": "",
      "methods": "",
      "sample": "",
      "manipulation": "",
      "measures": ""
    },
    "core_evidence": [
      {"title": "", "claim": ""},
      {"title": "", "claim": ""},
      {"title": "", "claim": ""},
      {"title": "", "claim": ""}
    ],
    "mechanism": "",
    "boundary_conditions": [],
    "generalizable_insight": ""
  },
  "claim_confidence": {
    "main_thesis": "supported",
    "mechanism": "hypothesis",
    "core_evidence": ["supported", "supported", "supported", "supported"]
  },
  "evidence_traceability": [
    {"claim_index": 1, "pages": [1], "support": "Short paraphrase of supporting evidence"},
    {"claim_index": 2, "pages": [2], "support": "Short paraphrase of supporting evidence"},
    {"claim_index": 3, "pages": [3], "support": "Short paraphrase of supporting evidence"},
    {"claim_index": 4, "pages": [4], "support": "Short paraphrase of supporting evidence"}
  ]
}

After generating scientific claims, compress them for infographic display.

Hard limits for infographic_claims, optimized for one A4 page:
- Thesis: max 24 words
- Why it matters: max 3 bullets, each max 11 words
- Study design fields: max 12 words each
- Core evidence: exactly 4 cards
- Each evidence card title: max 4 words
- Each evidence card claim: max 18 words
- Mechanism: max 28 words
- Boundary conditions: max 3 bullets, each max 12 words
- Generalizable insight: max 22 words

No paragraphs inside cards.
No semicolons.
No parenthetical overload.
Prefer short causal sentences.
Renderers will display only infographic_claims by default.
Every evidence card must have traceability to page markers when possible.
Do not fill space. Short, specific, and mechanistic beats complete but generic.
"""


FULL_CLAIMS_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": [
        "thesis",
        "why_it_matters",
        "study_design",
        "core_evidence",
        "mechanism",
        "boundary_conditions",
        "generalizable_insight",
    ],
    "properties": {
        "thesis": {"type": "string"},
        "why_it_matters": {"type": "array", "items": {"type": "string"}},
        "study_design": {
            "type": "object",
            "additionalProperties": False,
            "required": ["model_system", "methods", "sample", "manipulation", "measures"],
            "properties": {
                "model_system": {"type": "array", "items": {"type": "string"}},
                "methods": {"type": "array", "items": {"type": "string"}},
                "sample": {"type": "string"},
                "manipulation": {"type": "string"},
                "measures": {"type": "array", "items": {"type": "string"}},
            },
        },
        "core_evidence": {"type": "array", "items": {"type": "string"}},
        "mechanism": {"type": "string"},
        "boundary_conditions": {"type": "array", "items": {"type": "string"}},
        "generalizable_insight": {"type": "string"},
    },
}


INFOGRAPHIC_CLAIMS_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": [
        "thesis",
        "why_it_matters",
        "study_design",
        "core_evidence",
        "mechanism",
        "boundary_conditions",
        "generalizable_insight",
    ],
    "properties": {
        "thesis": {"type": "string"},
        "why_it_matters": {"type": "array", "items": {"type": "string"}},
        "study_design": {
            "type": "object",
            "additionalProperties": False,
            "required": ["model_system", "methods", "sample", "manipulation", "measures"],
            "properties": {
                "model_system": {"type": "string"},
                "methods": {"type": "string"},
                "sample": {"type": "string"},
                "manipulation": {"type": "string"},
                "measures": {"type": "string"},
            },
        },
        "core_evidence": {
            "type": "array",
            "items": {
                "type": "object",
                "additionalProperties": False,
                "required": ["title", "claim"],
                "properties": {
                    "title": {"type": "string"},
                    "claim": {"type": "string"},
                },
            },
        },
        "mechanism": {"type": "string"},
        "boundary_conditions": {"type": "array", "items": {"type": "string"}},
        "generalizable_insight": {"type": "string"},
    },
}


CLAIMS_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": ["full_structured_claims", "infographic_claims", "claim_confidence", "evidence_traceability"],
    "properties": {
        "full_structured_claims": FULL_CLAIMS_SCHEMA,
        "infographic_claims": INFOGRAPHIC_CLAIMS_SCHEMA,
        "claim_confidence": {
            "type": "object",
            "additionalProperties": False,
            "required": ["main_thesis", "mechanism", "core_evidence"],
            "properties": {
                "main_thesis": {"type": "string", "enum": ["strong", "supported", "speculative"]},
                "mechanism": {"type": "string", "enum": ["strong", "supported", "hypothesis", "speculative"]},
                "core_evidence": {
                    "type": "array",
                    "items": {"type": "string", "enum": ["strong", "supported", "speculative"]},
                    "minItems": 4,
                    "maxItems": 4,
                },
            },
        },
        "evidence_traceability": {
            "type": "array",
            "items": {
                "type": "object",
                "additionalProperties": False,
                "required": ["claim_index", "pages", "support"],
                "properties": {
                    "claim_index": {"type": "integer", "minimum": 1, "maximum": 4},
                    "pages": {"type": "array", "items": {"type": "integer"}, "minItems": 1, "maxItems": 4},
                    "support": {"type": "string"},
                },
            },
            "minItems": 4,
            "maxItems": 4,
        },
    },
}

QUALITY_CHECK_PROMPT = """You are a scientific claim critic.

Check whether the infographic claims are paper-specific, mechanistic, and supported by the supplied page-marked PDF text.
Do not rewrite the whole brief. Return compact JSON.
Flag generic, unsupported, or overconfident statements.
"""

QUALITY_CHECK_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": ["status", "risk_level", "issues", "recommended_fixes"],
    "properties": {
        "status": {"type": "string", "enum": ["pass", "review", "fail"]},
        "risk_level": {"type": "string", "enum": ["low", "medium", "high"]},
        "issues": {"type": "array", "items": {"type": "string"}, "maxItems": 5},
        "recommended_fixes": {"type": "array", "items": {"type": "string"}, "maxItems": 5},
    },
}

REGENERATE_PROMPT = """You rewrite one section of a scientific infographic.

Use only the supplied structured claims and metadata.
Keep it concise, mechanistic, paper-specific, and non-generic.
Return JSON with exactly one field: value.
"""

REGENERATE_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": ["value"],
    "properties": {"value": {"type": "string"}},
}

SYNTHESIS_PROMPT = """You are a cross-paper synthesis engine.

Use only the supplied structured claims.
Do not invent new paper-specific evidence.
Find the shared mechanism, disagreements, boundary conditions, and a general research implication.

Return compressed JSON for infographic display.
No generic statements.
No empty fields.
"""

SYNTHESIS_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": ["synthesis_thesis", "shared_mechanisms", "contrasts", "boundary_conditions", "research_implication"],
    "properties": {
        "synthesis_thesis": {"type": "string"},
        "shared_mechanisms": {"type": "array", "items": {"type": "string"}},
        "contrasts": {"type": "array", "items": {"type": "string"}},
        "boundary_conditions": {"type": "array", "items": {"type": "string"}},
        "research_implication": {"type": "string"},
    },
}

TITLE_LOOKUP_PROMPT = """You resolve scientific article metadata from DOI identifiers.

Use web search when available.
Return the exact article title matching the DOI.
Return the journal name.
Find one commonly reported journal-level metric when available, preferring Journal Impact Factor, CiteScore, SJR, or journal quartile.
Also find compact DOI-level context signals useful for an infographic:
- article type
- open access or paywall status
- citation or attention signal
- data/code availability
- funding, trial registration, preprint, or notable external context when available
- 2 to 4 concise reuse hooks explaining what makes the article useful to inspect
Assign an interest tier from the metric only:
- very_high: top-tier metric, usually Q1 or unusually high field-adjusted value
- high: strong journal metric
- moderate: visible but not top-tier metric
- low: weak, unavailable, or unverified metric
Do not return PDF filenames, journal section labels, publisher IDs, or guessed titles.
If a field cannot be verified, return "not found" and low confidence.
"""

TITLE_LOOKUP_SCHEMA: dict[str, Any] = {
    "type": "object",
    "additionalProperties": False,
    "required": ["doi", "title", "journal", "confidence", "source_url", "journal_metric", "article_signals"],
    "properties": {
        "doi": {"type": "string"},
        "title": {"type": "string"},
        "journal": {"type": "string"},
        "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
        "source_url": {"type": "string"},
        "journal_metric": {
            "type": "object",
            "additionalProperties": False,
            "required": [
                "metric_name",
                "metric_value",
                "metric_year",
                "quartile",
                "interest_score",
                "interest_tier",
                "rationale",
                "source_url",
            ],
            "properties": {
                "metric_name": {"type": "string"},
                "metric_value": {"type": "string"},
                "metric_year": {"type": "string"},
                "quartile": {"type": "string"},
                "interest_score": {"type": "integer", "minimum": 0, "maximum": 100},
                "interest_tier": {"type": "string", "enum": ["low", "moderate", "high", "very_high"]},
                "rationale": {"type": "string"},
                "source_url": {"type": "string"},
            },
        },
        "article_signals": {
            "type": "object",
            "additionalProperties": False,
            "required": [
                "article_type",
                "access_status",
                "citation_signal",
                "data_code_signal",
                "external_context",
                "reuse_hooks",
                "source_urls",
            ],
            "properties": {
                "article_type": {"type": "string"},
                "access_status": {"type": "string"},
                "citation_signal": {"type": "string"},
                "data_code_signal": {"type": "string"},
                "external_context": {"type": "string"},
                "reuse_hooks": {"type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 4},
                "source_urls": {"type": "array", "items": {"type": "string"}, "maxItems": 4},
            },
        },
    },
}


ROOT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BATCH_STORE: dict[str, dict[str, Any]] = {}
DEFAULT_MODEL_OPTIONS = ["gpt-5.2", "gpt-5.1", "gpt-5", "gpt-4.1", "gpt-4.1-mini", "o4-mini"]

app = FastAPI(title="PaperBrief API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)
app.mount("/static", StaticFiles(directory=ROOT_DIR), name="static")


@app.get("/")
def index() -> FileResponse:
    return FileResponse(os.path.join(ROOT_DIR, "index.html"))


@app.get("/app.js")
def app_js() -> FileResponse:
    return FileResponse(os.path.join(ROOT_DIR, "app.js"))


@app.get("/styles.css")
def styles_css() -> FileResponse:
    return FileResponse(os.path.join(ROOT_DIR, "styles.css"))


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


def clean_pdf_line(line: str) -> str:
    clean = re.sub(r"\s+", " ", line).strip()
    return clean.strip(" -•·")


def extract_doi(text: str) -> str:
    match = re.search(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+\b", text, re.IGNORECASE)
    if not match:
        return ""
    return match.group(0).rstrip(".,;:)])}>").lower()


def looks_like_internal_title(title: str, filename: str) -> bool:
    value = clean_pdf_line(title).lower()
    if not value:
        return True
    file_stem = os.path.splitext(os.path.basename(filename))[0].lower()
    if value == file_stem or file_stem in value:
        return True
    return bool(re.search(r"[_/\\]|\.pdf$|\+\+|\.{2,}|\b\d+\.\.\d+\b", value))


def title_candidate_is_noise(candidate: str, filename: str) -> bool:
    lowered = candidate.lower().rstrip(":")
    noise_patterns = [
        r"^research article$",
        r"^original article$",
        r"^review article$",
        r"^open access$",
        r"^article$",
        r"^translational neuroscience$",
        r"^de gruyter$",
        r"^received\b",
        r"^accepted\b",
        r"^published\b",
        r"^correspond",
        r"^copyright\b",
        r"^doi\b",
        r"^https?://",
        r"@",
    ]
    if looks_like_internal_title(candidate, filename):
        return True
    if any(re.search(pattern, lowered) for pattern in noise_patterns):
        return True
    if re.search(r"\b(university|department|faculty|institute|hospital|clinic)\b", lowered):
        return True
    if re.search(r"\b(author|license|creative commons|citation)\b", lowered):
        return True
    if lowered.count(",") >= 3:
        return True
    return False


def infer_layout_title(first_page: fitz.Page, filename: str) -> str:
    page_dict = first_page.get_text("dict")
    lines: list[dict[str, Any]] = []
    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [span for span in line.get("spans", []) if span.get("text", "").strip()]
            text = clean_pdf_line(" ".join(span["text"] for span in spans))
            if not text or title_candidate_is_noise(text, filename):
                continue
            words = text.split()
            if len(words) > 18:
                continue
            max_size = max(float(span.get("size", 0)) for span in spans)
            y0 = float(line.get("bbox", [0, 9999, 0, 9999])[1])
            x0 = float(line.get("bbox", [9999, 0, 9999, 0])[0])
            if y0 > first_page.rect.height * 0.45:
                continue
            lines.append({"text": text, "size": max_size, "y": y0, "x": x0})

    if not lines:
        return ""

    max_size = max(line["size"] for line in lines)
    title_lines = [
        line for line in lines
        if line["size"] >= max_size - 0.5 and line["y"] <= first_page.rect.height * 0.35
    ]
    title_lines.sort(key=lambda item: (item["y"], item["x"]))
    title = clean_pdf_line(" ".join(line["text"] for line in title_lines))
    if 4 <= len(title.split()) <= 30 and not title_candidate_is_noise(title, filename):
        return title.rstrip(".")
    return ""


def infer_article_title(first_page_text: str, pdf_title: str, filename: str) -> str:
    if pdf_title and not looks_like_internal_title(pdf_title, filename):
        return clean_pdf_line(pdf_title)

    stop_markers = {
        "abstract",
        "keywords",
        "introduction",
        "references",
        "methods",
        "materials and methods",
    }
    lines: list[str] = []
    for raw_line in first_page_text.splitlines():
        line = clean_pdf_line(raw_line)
        if not line:
            continue
        lowered = line.lower().rstrip(":")
        if lowered in stop_markers:
            break
        if title_candidate_is_noise(line, filename):
            continue
        if re.search(r"\b\d{4}\b", lowered) and len(line.split()) <= 5:
            continue
        lines.append(line)
        if len(lines) >= 40:
            break

    candidates: list[tuple[int, str]] = []
    for start in range(len(lines)):
        for length in range(1, 4):
            chunk = lines[start:start + length]
            if len(chunk) != length:
                continue
            candidate = clean_pdf_line(" ".join(chunk))
            words = candidate.split()
            if not 5 <= len(words) <= 28:
                continue
            lowered = candidate.lower()
            if title_candidate_is_noise(candidate, filename):
                continue

            score = 100 - abs(14 - len(words)) * 2 - start
            if any(char in candidate for char in [":", "?", "-"]):
                score += 4
            if candidate.endswith("."):
                score -= 8
            candidates.append((score, candidate.rstrip(".")))

    if candidates:
        return max(candidates, key=lambda item: item[0])[1]

    if pdf_title and not looks_like_internal_title(pdf_title, filename):
        return clean_pdf_line(pdf_title)
    return ""


def extract_pdf_text(pdf_bytes: bytes, filename: str = "") -> dict[str, Any]:
    try:
        document = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not open PDF: {exc}") from exc

    pages: list[dict[str, Any]] = []
    for page_index, page in enumerate(document, start=1):
        text = page.get_text("text").strip()
        pages.append({"page": page_index, "text": text})

    raw_text = "\n\n".join(f"[PAGE {page['page']}]\n{page['text']}" for page in pages if page["text"])
    pdf_metadata = document.metadata or {}
    pdf_title = pdf_metadata.get("title") or ""
    first_page_text = pages[0]["text"] if pages else ""
    layout_title = infer_layout_title(document[0], filename) if document.page_count else ""
    fallback_title = layout_title or infer_article_title(first_page_text, pdf_title, filename)
    metadata = {
        "title": fallback_title,
        "fallback_title": fallback_title,
        "title_source": "pdf_fallback" if fallback_title else "not_detected",
        "doi": extract_doi(raw_text),
        "journal": "",
        "journal_metric": normalize_journal_metric(None),
        "article_signals": normalize_article_signals(None),
        "pdf_metadata_title": pdf_title,
        "author": pdf_metadata.get("author") or "",
        "page_count": document.page_count,
        "char_count": len(raw_text),
    }
    snapshot = render_first_page_snapshot(document)
    document.close()

    if not raw_text.strip():
        raise HTTPException(status_code=400, detail="No extractable text found. The PDF may be scanned or encrypted.")

    return {"text": raw_text, "metadata": metadata, "snapshot": snapshot}


def render_first_page_snapshot(document: fitz.Document) -> dict[str, Any] | None:
    if document.page_count < 1:
        return None
    try:
        page = document[0]
        zoom = min(1.8, 900 / max(page.rect.width, 1))
        pixmap = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        png_bytes = pixmap.tobytes("png")
        return {
            "page": 1,
            "width": pixmap.width,
            "height": pixmap.height,
            "data_url": f"data:image/png;base64,{base64.b64encode(png_bytes).decode('ascii')}",
        }
    except Exception:
        return None


TEMPLATE_PRESETS = {
    "scientific_claims": "Default: extract compressed causal scientific claims.",
    "clinical_relevance": "Emphasize clinical relevance, diagnosis, intervention, patient stratification, and translational limits.",
    "journal_club": "Emphasize study design, interpretation risks, discussion questions, and what a journal club should debate.",
    "grant_research_idea": "Emphasize research gaps, next experiments, and grant-worthy mechanistic questions.",
    "ai_systems_analogy": "Emphasize abstract mechanisms and analogies to AI, neural networks, and complex systems.",
}


def build_user_input(
    raw_text: str,
    metadata: dict[str, Any],
    filename: str,
    template_preset: str = "scientific_claims",
    custom_instructions: str = "",
) -> str:
    max_chars = 120_000
    clipped_text = raw_text[:max_chars]
    clipping_note = ""
    if len(raw_text) > max_chars:
        clipping_note = f"\n\n[NOTE: PDF text was clipped from {len(raw_text)} to {max_chars} characters.]"

    preset = TEMPLATE_PRESETS.get(template_preset, TEMPLATE_PRESETS["scientific_claims"])
    custom = clean_pdf_line(custom_instructions)[:1200]

    return f"""Extract structured scientific claims from this paper.

Additional instructions:
- prioritize mechanism over description
- convert descriptive findings into causal statements
- produce a strong one-line thesis
- include a generalizable insight
- do not copy paper sentences verbatim
- reject generic background claims
- produce claim confidence and evidence traceability
- use page markers like [PAGE 3] for evidence_traceability pages
- template preset: {preset}
- user instruction profile: {custom or "none"}

PDF metadata:
- filename: {filename}
- title: {metadata.get("title") or "not detected"}
- doi: {metadata.get("doi") or "not detected"}
- journal: {metadata.get("journal") or "not detected"}
- author: {metadata.get("author") or "not detected"}
- pages: {metadata.get("page_count")}

Input:
{clipped_text}{clipping_note}
"""


def resolve_api_key(api_key: str | None) -> str:
    key = (api_key or os.getenv("OPENAI_API_KEY") or "").strip()
    if not key:
        raise HTTPException(status_code=400, detail="Provide an API key or set OPENAI_API_KEY.")
    if not key.startswith("sk-"):
        raise HTTPException(status_code=401, detail="OpenAI API key looks invalid. It should start with sk-.")
    return key


def is_usable_generation_model(model_id: str) -> bool:
    if not model_id.startswith(("gpt-", "o")):
        return False
    excluded = [
        "audio",
        "chat",
        "codex",
        "embedding",
        "image",
        "pro",
        "realtime",
        "search",
        "speech",
        "tts",
        "transcribe",
        "whisper",
    ]
    if any(token in model_id for token in excluded):
        return False
    if model_id.startswith("gpt-5"):
        return True
    if model_id.startswith(("gpt-4.1", "o4")):
        return True
    return False


def sort_model_ids(model_id: str) -> tuple[int, str]:
    if model_id in DEFAULT_MODEL_OPTIONS:
        return (DEFAULT_MODEL_OPTIONS.index(model_id), model_id)
    if model_id.startswith("gpt-5"):
        return (20, model_id)
    if model_id.startswith("gpt-4"):
        return (30, model_id)
    if model_id.startswith("o"):
        return (40, model_id)
    return (90, model_id)


@app.post("/api/models")
def list_openai_models(api_key: str | None = Body(default=None, embed=True)) -> dict[str, Any]:
    key = resolve_api_key(api_key)
    client = OpenAI(api_key=key)
    try:
        response = client.models.list()
    except AuthenticationError as exc:
        raise HTTPException(status_code=401, detail="OpenAI rejected the API key. Paste a valid key or set OPENAI_API_KEY.") from exc
    except OpenAIError as exc:
        raise HTTPException(status_code=502, detail=f"Could not load OpenAI models: {exc}") from exc

    model_ids = sorted(
        {model.id for model in response.data if is_usable_generation_model(model.id)},
        key=sort_model_ids,
    )
    if not model_ids:
        model_ids = DEFAULT_MODEL_OPTIONS
    return {"models": model_ids}


def call_openai(api_key: str, model: str, user_input: str) -> dict[str, Any]:
    client = OpenAI(api_key=api_key)
    try:
        response = client.responses.create(
            model=model,
            instructions=SYSTEM_PROMPT,
            input=user_input,
            text={
                "format": {
                    "type": "json_schema",
                    "name": "structured_scientific_claims",
                    "strict": True,
                    "schema": CLAIMS_SCHEMA,
                }
            },
        )
    except AuthenticationError as exc:
        raise HTTPException(status_code=401, detail="OpenAI rejected the API key. Paste a valid key or set OPENAI_API_KEY.") from exc
    except OpenAIError as exc:
        raise HTTPException(status_code=502, detail=f"OpenAI API request failed: {exc}") from exc

    try:
        claims = json.loads(response.output_text)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Model did not return valid JSON: {exc}") from exc
    normalize_confidence_and_traceability(claims)
    normalization_notes = normalize_infographic_claims(claims)
    validate_claims(claims)
    claims["_normalization_notes"] = normalization_notes
    return claims


def call_openai_quality_check(api_key: str, model: str, raw_text: str, claims: dict[str, Any]) -> dict[str, Any]:
    client = OpenAI(api_key=api_key)
    quality_input = {
        "pdf_text": raw_text[:80_000],
        "claims": claims,
    }
    try:
        response = client.responses.create(
            model=model,
            instructions=QUALITY_CHECK_PROMPT,
            input=json.dumps(quality_input),
            text={
                "format": {
                    "type": "json_schema",
                    "name": "claim_quality_check",
                    "strict": True,
                    "schema": QUALITY_CHECK_SCHEMA,
                }
            },
        )
        return json.loads(response.output_text)
    except (OpenAIError, Exception):
        return {
            "status": "review",
            "risk_level": "medium",
            "issues": ["Automated quality check unavailable"],
            "recommended_fixes": ["Manually verify claims against the PDF"],
        }


def normalize_confidence_and_traceability(claims: dict[str, Any]) -> None:
    confidence = claims.get("claim_confidence")
    if not isinstance(confidence, dict):
        confidence = {}
    core_confidence = confidence.get("core_evidence")
    if not isinstance(core_confidence, list):
        core_confidence = []
    allowed_core = {"strong", "supported", "speculative"}
    allowed_mechanism = {"strong", "supported", "hypothesis", "speculative"}
    claims["claim_confidence"] = {
        "main_thesis": confidence.get("main_thesis") if confidence.get("main_thesis") in allowed_core else "supported",
        "mechanism": confidence.get("mechanism") if confidence.get("mechanism") in allowed_mechanism else "hypothesis",
        "core_evidence": [
            item if item in allowed_core else "supported"
            for item in (core_confidence + ["supported", "supported", "supported", "supported"])[:4]
        ],
    }

    traceability = claims.get("evidence_traceability")
    if not isinstance(traceability, list):
        traceability = []
    normalized_trace = []
    for index in range(1, 5):
        item = next((entry for entry in traceability if isinstance(entry, dict) and entry.get("claim_index") == index), {})
        pages = item.get("pages") if isinstance(item, dict) else []
        if not isinstance(pages, list) or not pages:
            pages = [1]
        normalized_trace.append({
            "claim_index": index,
            "pages": [int(page) for page in pages[:4] if str(page).isdigit()] or [1],
            "support": clean_pdf_line(str(item.get("support") or "Support should be verified against the PDF"))[:180],
        })
    claims["evidence_traceability"] = normalized_trace


def call_openai_title_lookup(api_key: str, model: str, metadata: dict[str, Any], filename: str) -> dict[str, Any] | None:
    doi = (metadata.get("doi") or "").strip()
    if not doi:
        return None

    client = OpenAI(api_key=api_key)
    lookup_input = {
        "doi": doi,
        "filename": filename,
        "pdf_fallback_title": metadata.get("fallback_title") or metadata.get("title") or "",
        "pdf_metadata_title": metadata.get("pdf_metadata_title") or "",
    }
    try:
        response = client.responses.create(
            model=model,
            instructions=TITLE_LOOKUP_PROMPT,
            input=json.dumps(lookup_input),
            tools=[{"type": "web_search_preview"}],
            text={
                "format": {
                    "type": "json_schema",
                    "name": "doi_title_lookup",
                    "strict": True,
                    "schema": TITLE_LOOKUP_SCHEMA,
                }
            },
        )
    except AuthenticationError as exc:
        raise HTTPException(status_code=401, detail="OpenAI rejected the API key. Paste a valid key or set OPENAI_API_KEY.") from exc
    except OpenAIError:
        return None

    try:
        result = json.loads(response.output_text)
    except Exception:
        return None

    title = clean_pdf_line(result.get("title", ""))
    result["title"] = title.rstrip(".")
    result["journal"] = clean_pdf_line(result.get("journal", ""))
    return result


def resolve_title_with_openai(api_key: str, model: str, parsed: dict[str, Any], filename: str) -> None:
    metadata = parsed["metadata"]
    if not metadata.get("doi"):
        return

    lookup = call_openai_title_lookup(api_key, model, metadata, filename)
    if not lookup:
        metadata["title_source"] = "pdf_fallback_after_doi_lookup"
        return

    title = lookup.get("title", "")
    if title and lookup.get("confidence") != "low" and not title_candidate_is_noise(title, filename):
        metadata["title"] = title
        metadata["title_source"] = "openai_doi_lookup"
        metadata["title_lookup_confidence"] = lookup.get("confidence", "")
        metadata["title_lookup_source_url"] = lookup.get("source_url", "")
    else:
        metadata["title_source"] = "pdf_fallback_after_doi_lookup"
    metadata["journal"] = lookup.get("journal", "")
    metadata["journal_metric"] = normalize_journal_metric(lookup.get("journal_metric"))
    metadata["article_signals"] = normalize_article_signals(lookup.get("article_signals"))


def normalize_journal_metric(metric: Any) -> dict[str, Any]:
    if not isinstance(metric, dict):
        metric = {}
    tier = metric.get("interest_tier") if metric.get("interest_tier") in {"low", "moderate", "high", "very_high"} else "low"
    try:
        score = int(metric.get("interest_score", 0))
    except (TypeError, ValueError):
        score = 0
    return {
        "metric_name": clean_pdf_line(str(metric.get("metric_name") or "not found")),
        "metric_value": clean_pdf_line(str(metric.get("metric_value") or "not found")),
        "metric_year": clean_pdf_line(str(metric.get("metric_year") or "not found")),
        "quartile": clean_pdf_line(str(metric.get("quartile") or "not found")),
        "interest_score": max(0, min(100, score)),
        "interest_tier": tier,
        "rationale": clean_pdf_line(str(metric.get("rationale") or "Metric unavailable or unverified")),
        "source_url": clean_pdf_line(str(metric.get("source_url") or "")),
    }


def normalize_article_signals(signals: Any) -> dict[str, Any]:
    if not isinstance(signals, dict):
        signals = {}
    hooks = signals.get("reuse_hooks")
    if not isinstance(hooks, list):
        hooks = []
    clean_hooks = [
        clean_pdf_line(str(item))
        for item in hooks
        if clean_pdf_line(str(item)) and clean_pdf_line(str(item)).lower() != "not found"
    ][:4]

    source_urls = signals.get("source_urls")
    if not isinstance(source_urls, list):
        source_urls = []
    clean_urls = [
        clean_pdf_line(str(url))
        for url in source_urls
        if clean_pdf_line(str(url)) and clean_pdf_line(str(url)).lower() != "not found"
    ][:4]

    return {
        "article_type": clean_pdf_line(str(signals.get("article_type") or "not found")),
        "access_status": clean_pdf_line(str(signals.get("access_status") or "not found")),
        "citation_signal": clean_pdf_line(str(signals.get("citation_signal") or "not found")),
        "data_code_signal": clean_pdf_line(str(signals.get("data_code_signal") or "not found")),
        "external_context": clean_pdf_line(str(signals.get("external_context") or "not found")),
        "reuse_hooks": clean_hooks,
        "source_urls": clean_urls,
    }


def call_openai_synthesis(api_key: str, model: str, papers: list[dict[str, Any]]) -> dict[str, Any]:
    client = OpenAI(api_key=api_key)
    synthesis_input = {
        "task": "Cross-paper synthesis",
        "papers": [
            {
                "filename": paper["filename"],
                "metadata": paper["metadata"],
                "claims": paper["claims"]["full_structured_claims"],
                "infographic": paper["claims"]["infographic_claims"],
            }
            for paper in papers
        ],
    }
    try:
        response = client.responses.create(
            model=model,
            instructions=SYNTHESIS_PROMPT,
            input=json.dumps(synthesis_input),
            text={
                "format": {
                    "type": "json_schema",
                    "name": "cross_paper_synthesis",
                    "strict": True,
                    "schema": SYNTHESIS_SCHEMA,
                }
            },
        )
    except AuthenticationError as exc:
        raise HTTPException(status_code=401, detail="OpenAI rejected the API key. Paste a valid key or set OPENAI_API_KEY.") from exc
    except OpenAIError as exc:
        raise HTTPException(status_code=502, detail=f"OpenAI synthesis request failed: {exc}") from exc
    try:
        return json.loads(response.output_text)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Synthesis did not return valid JSON: {exc}") from exc


def word_count(text: str) -> int:
    return len(str(text).replace("→", " ").split())


def strip_overload(text: str) -> str:
    clean = str(text).replace(";", ",").strip()
    while "(" in clean and ")" in clean:
        start = clean.find("(")
        end = clean.find(")", start)
        if end == -1:
            break
        clean = f"{clean[:start].rstrip()} {clean[end + 1:].lstrip()}".strip()
    return " ".join(clean.split())


def limit_words(text: str, limit: int) -> str:
    words = strip_overload(text).split()
    if len(words) <= limit:
        return " ".join(words)
    return " ".join(words[:limit]).rstrip(" ,.")


def normalize_infographic_claims(claims: dict[str, Any]) -> list[str]:
    infographic = claims.get("infographic_claims")
    if not isinstance(infographic, dict):
        return []

    notes: list[str] = []

    def normalize_text(label: str, value: Any, limit: int) -> str:
        original = strip_overload(str(value or ""))
        normalized = limit_words(original, limit)
        if normalized != original:
            notes.append(f"{label} compressed")
        return normalized or "Not specified"

    infographic["thesis"] = normalize_text("thesis", infographic.get("thesis", ""), 24)
    infographic["why_it_matters"] = [
        normalize_text(f"why_it_matters[{index + 1}]", item, 11)
        for index, item in enumerate((infographic.get("why_it_matters") or [])[:3])
    ] or ["Paper changes the mechanistic interpretation of the problem"]

    design = infographic.get("study_design") or {}
    infographic["study_design"] = {
        "model_system": normalize_text("study_design.model_system", design.get("model_system", ""), 12),
        "methods": normalize_text("study_design.methods", design.get("methods", ""), 12),
        "sample": normalize_text("study_design.sample", design.get("sample", ""), 12),
        "manipulation": normalize_text("study_design.manipulation", design.get("manipulation", ""), 12),
        "measures": normalize_text("study_design.measures", design.get("measures", ""), 12),
    }

    evidence = list(infographic.get("core_evidence") or [])
    while len(evidence) < 4:
        evidence.append({"title": f"Claim {len(evidence) + 1}", "claim": "Mechanistic claim requires model retry"})
    infographic["core_evidence"] = [
        {
            "title": normalize_text(f"core_evidence[{index + 1}].title", card.get("title", ""), 4),
            "claim": normalize_text(f"core_evidence[{index + 1}].claim", card.get("claim", ""), 18),
        }
        for index, card in enumerate(evidence[:4])
    ]

    infographic["mechanism"] = normalize_text("mechanism", infographic.get("mechanism", ""), 28)
    infographic["boundary_conditions"] = [
        normalize_text(f"boundary_conditions[{index + 1}]", item, 12)
        for index, item in enumerate((infographic.get("boundary_conditions") or [])[:3])
    ] or ["Boundary conditions require closer reading"]
    infographic["generalizable_insight"] = normalize_text(
        "generalizable_insight", infographic.get("generalizable_insight", ""), 22
    )
    return notes


def validate_claims(claims: dict[str, Any]) -> None:
    def empty(value: Any) -> bool:
        if isinstance(value, str):
            return not value.strip()
        if isinstance(value, list):
            return len(value) == 0 or any(empty(item) for item in value)
        if isinstance(value, dict):
            return any(empty(item) for item in value.values())
        return value is None

    if empty(claims):
        raise HTTPException(status_code=502, detail="Model returned empty fields; retry generation.")

    infographic = claims.get("infographic_claims", {})
    errors: list[str] = []

    def check_words(label: str, text: str, limit: int) -> None:
        if word_count(text) > limit:
            errors.append(f"{label} exceeds {limit} words")
        if ";" in str(text):
            errors.append(f"{label} contains a semicolon")

    check_words("thesis", infographic.get("thesis", ""), 24)
    for index, item in enumerate(infographic.get("why_it_matters", []), start=1):
      check_words(f"why_it_matters[{index}]", item, 11)
    if len(infographic.get("why_it_matters", [])) > 3:
        errors.append("why_it_matters exceeds 3 bullets")

    design = infographic.get("study_design", {})
    for key in ["model_system", "methods", "sample", "manipulation", "measures"]:
        check_words(f"study_design.{key}", design.get(key, ""), 12)

    evidence = infographic.get("core_evidence", [])
    if len(evidence) != 4:
        errors.append("core_evidence must contain exactly 4 cards")
    for index, card in enumerate(evidence, start=1):
        check_words(f"core_evidence[{index}].title", card.get("title", ""), 4)
        check_words(f"core_evidence[{index}].claim", card.get("claim", ""), 18)

    check_words("mechanism", infographic.get("mechanism", ""), 28)
    for index, item in enumerate(infographic.get("boundary_conditions", []), start=1):
        check_words(f"boundary_conditions[{index}]", item, 12)
    if len(infographic.get("boundary_conditions", [])) > 3:
        errors.append("boundary_conditions exceeds 3 bullets")
    check_words("generalizable_insight", infographic.get("generalizable_insight", ""), 22)

    if errors:
        raise HTTPException(status_code=502, detail=f"Model returned non-compliant infographic_claims: {', '.join(errors)}")


@app.post("/api/generate")
async def generate_claims(
    pdf: UploadFile = File(...),
    api_key: str | None = Form(default=None),
    model: str = Form(default="gpt-5.2"),
    template_preset: str = Form(default="scientific_claims"),
    custom_instructions: str = Form(default=""),
) -> dict[str, Any]:
    key = resolve_api_key(api_key)

    if not pdf.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Upload a PDF file.")

    pdf_bytes = await pdf.read()
    parsed = extract_pdf_text(pdf_bytes, pdf.filename)
    resolve_title_with_openai(key, model, parsed, pdf.filename)
    user_input = build_user_input(parsed["text"], parsed["metadata"], pdf.filename, template_preset, custom_instructions)
    claims = call_openai(key, model, user_input)
    claims["quality_check"] = call_openai_quality_check(key, model, parsed["text"], claims)

    return {
        "filename": pdf.filename,
        "metadata": parsed["metadata"],
        "snapshot": parsed.get("snapshot"),
        "model": model,
        "claims": claims,
    }


@app.post("/api/generate-batch")
async def generate_batch(
    pdfs: list[UploadFile] = File(...),
    api_key: str | None = Form(default=None),
    model: str = Form(default="gpt-5.2"),
    synthesis_mode: str = Form(default="separate"),
    template_preset: str = Form(default="scientific_claims"),
    custom_instructions: str = Form(default=""),
) -> dict[str, Any]:
    key = resolve_api_key(api_key)
    if not pdfs:
        raise HTTPException(status_code=400, detail="Upload at least one PDF.")
    normalized_mode = synthesis_mode if synthesis_mode in {"separate", "synthesis"} else "separate"

    papers: list[dict[str, Any]] = []
    for pdf in pdfs:
        if not pdf.filename.lower().endswith(".pdf"):
            raise HTTPException(status_code=400, detail=f"{pdf.filename} is not a PDF.")
        parsed = extract_pdf_text(await pdf.read(), pdf.filename)
        resolve_title_with_openai(key, model, parsed, pdf.filename)
        user_input = build_user_input(parsed["text"], parsed["metadata"], pdf.filename, template_preset, custom_instructions)
        claims = call_openai(key, model, user_input)
        claims["quality_check"] = call_openai_quality_check(key, model, parsed["text"], claims)
        papers.append({
            "filename": pdf.filename,
            "metadata": parsed["metadata"],
            "snapshot": parsed.get("snapshot"),
            "model": model,
            "claims": claims,
        })

    synthesis = None
    if len(papers) > 1 and normalized_mode == "synthesis":
        synthesis = call_openai_synthesis(key, model, papers)

    batch_id = str(uuid.uuid4())
    payload = {
        "batch_id": batch_id,
        "model": model,
        "synthesis_mode": normalized_mode,
        "papers": papers,
        "synthesis": synthesis,
        "template_preset": template_preset,
    }
    BATCH_STORE[batch_id] = payload
    return payload


@app.post("/api/regenerate-section")
def regenerate_section(payload: dict[str, Any] = Body(...)) -> dict[str, Any]:
    key = resolve_api_key(payload.get("api_key"))
    model = payload.get("model") or "gpt-5.2"
    section = payload.get("section") or "thesis"
    paper = payload.get("paper") or {}
    style = clean_pdf_line(str(payload.get("style") or ""))[:500]
    client = OpenAI(api_key=key)
    request = {
        "section": section,
        "style": style or "Improve clarity and specificity",
        "metadata": paper.get("metadata", {}),
        "claims": paper.get("claims", {}),
    }
    try:
        response = client.responses.create(
            model=model,
            instructions=REGENERATE_PROMPT,
            input=json.dumps(request),
            text={
                "format": {
                    "type": "json_schema",
                    "name": "regenerated_section",
                    "strict": True,
                    "schema": REGENERATE_SCHEMA,
                }
            },
        )
        value = json.loads(response.output_text).get("value", "")
    except AuthenticationError as exc:
        raise HTTPException(status_code=401, detail="OpenAI rejected the API key. Paste a valid key or set OPENAI_API_KEY.") from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Section regeneration failed: {exc}") from exc
    return {"section": section, "value": limit_words(value, 35)}


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=RGBColor(29, 36, 35)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(size)
    return box


def build_pptx(batch: dict[str, Any]) -> str:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]

    cover = deck.slides.add_slide(blank)
    add_textbox(cover, Inches(0.7), Inches(0.65), Inches(11.8), Inches(0.4), "PaperBrief", 14, True, RGBColor(216, 111, 87))
    add_textbox(cover, Inches(0.7), Inches(1.35), Inches(11.8), Inches(1.0), "Scientific Claims Deck", 36, True)
    add_textbox(cover, Inches(0.7), Inches(2.35), Inches(10.5), Inches(0.6), f"{len(batch['papers'])} paper(s) · generated from OpenAI structured claims", 18)

    if batch.get("synthesis"):
        synthesis = batch["synthesis"]
        slide = deck.slides.add_slide(blank)
        add_textbox(slide, Inches(0.7), Inches(0.55), Inches(11.8), Inches(0.5), "Cross-paper synthesis", 30, True)
        add_textbox(slide, Inches(0.7), Inches(1.25), Inches(11.8), Inches(0.8), synthesis["synthesis_thesis"], 20, True)
        add_textbox(slide, Inches(0.7), Inches(2.25), Inches(3.5), Inches(0.3), "Shared mechanisms", 13, True, RGBColor(31, 122, 120))
        add_bullets(slide, Inches(0.7), Inches(2.6), Inches(3.7), Inches(2.3), synthesis["shared_mechanisms"][:4], 13)
        add_textbox(slide, Inches(4.8), Inches(2.25), Inches(3.5), Inches(0.3), "Contrasts", 13, True, RGBColor(31, 122, 120))
        add_bullets(slide, Inches(4.8), Inches(2.6), Inches(3.7), Inches(2.3), synthesis["contrasts"][:4], 13)
        add_textbox(slide, Inches(8.9), Inches(2.25), Inches(3.5), Inches(0.3), "Research implication", 13, True, RGBColor(31, 122, 120))
        add_textbox(slide, Inches(8.9), Inches(2.6), Inches(3.7), Inches(1.4), synthesis["research_implication"], 15, True)

    for paper in batch["papers"]:
        claims = paper["claims"]["infographic_claims"]
        slide = deck.slides.add_slide(blank)
        metadata = paper.get("metadata") or {}
        title = metadata.get("title") or claims.get("title") or paper.get("filename", "Uploaded paper")
        add_textbox(slide, Inches(0.7), Inches(0.45), Inches(10.8), Inches(0.55), title[:120], 24, True)
        add_textbox(slide, Inches(0.7), Inches(1.1), Inches(11.7), Inches(0.6), claims["thesis"], 18, True, RGBColor(31, 122, 120))
        add_textbox(slide, Inches(0.7), Inches(2.0), Inches(3.6), Inches(0.3), "Why it matters", 12, True, RGBColor(216, 111, 87))
        add_bullets(slide, Inches(0.7), Inches(2.35), Inches(3.7), Inches(1.6), claims["why_it_matters"], 12)
        add_textbox(slide, Inches(4.8), Inches(2.0), Inches(3.6), Inches(0.3), "Core evidence", 12, True, RGBColor(216, 111, 87))
        evidence_lines = [f"{card['title']}: {card['claim']}" for card in claims["core_evidence"]]
        add_bullets(slide, Inches(4.8), Inches(2.35), Inches(4.0), Inches(2.3), evidence_lines, 12)
        add_textbox(slide, Inches(9.2), Inches(2.0), Inches(3.2), Inches(0.3), "Mechanism", 12, True, RGBColor(216, 111, 87))
        add_textbox(slide, Inches(9.2), Inches(2.35), Inches(3.3), Inches(1.2), claims["mechanism"], 13, True)
        add_textbox(slide, Inches(0.7), Inches(5.45), Inches(11.8), Inches(0.5), claims["generalizable_insight"], 18, True)

    export_dir = os.path.join(os.getcwd(), "exports")
    os.makedirs(export_dir, exist_ok=True)
    output = os.path.join(export_dir, f"paperbrief-{batch.get('batch_id', uuid.uuid4().hex)}.pptx")
    deck.save(output)
    return output


@app.get("/api/export-pptx/{batch_id}")
def export_pptx(batch_id: str) -> FileResponse:
    batch = BATCH_STORE.get(batch_id)
    if not batch:
        raise HTTPException(status_code=404, detail="Batch not found. Generate a batch first.")
    path = build_pptx(batch)
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="paperbrief-claims.pptx",
    )


@app.post("/api/export-pptx")
def export_pptx_payload(batch: dict[str, Any] = Body(...)) -> FileResponse:
    if not batch.get("papers"):
        raise HTTPException(status_code=400, detail="No papers to export.")
    path = build_pptx(batch)
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="paperbrief-claims.pptx",
    )
